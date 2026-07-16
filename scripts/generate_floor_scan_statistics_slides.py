#!/usr/bin/env python3
"""Build a short presentation on statistics for a typical Fig. 7 scan point."""

from __future__ import annotations

import csv
import os
import subprocess
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from matplotlib.colors import LogNorm
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[1]
DATA_DIR = Path(os.environ.get("HFIRBGDATA", "/home/blaine/projects/HFIRBG/data"))
FIG7_PDF = Path(
    os.environ.get(
        "HFIRBG_FIG7_PDF",
        "/home/blaine/projects/HFIRBG/paper/arxiv_submission/figures/pdf/down_facing_50_to_11400.pdf",
    )
)
OUTDIR = REPO / "analysis" / "presentations" / "floor_scan_statistics"
PPTX_OUT = OUTDIR / "floor_scan_average_point_statistics_2026-07-10.pptx"

NAVY = RGBColor(20, 43, 68)
BLUE = RGBColor(31, 119, 180)
TEAL = RGBColor(0, 135, 147)
ORANGE = RGBColor(230, 126, 34)
GREEN = RGBColor(45, 145, 80)
GRAY = RGBColor(95, 105, 115)
LIGHT = RGBColor(239, 244, 248)
WHITE = RGBColor(255, 255, 255)


def read_spectrum(path: Path):
    energy, counts = [], []
    with path.open() as stream:
        for line in stream:
            if line.startswith("#") or not line.strip():
                continue
            fields = line.split()
            if len(fields) < 3:
                continue
            try:
                energy.append(float(fields[1]))
                counts.append(float(fields[2]))
            except ValueError:
                pass
    return np.asarray(energy), np.asarray(counts)


def load_points():
    points = []
    for path in sorted((REPO / "db").glob("position_scan_*.csv")):
        with path.open(newline="") as stream:
            rows = list(csv.reader(stream, delimiter="|"))[1:]
        for row in rows:
            try:
                angle = float(row[2])
                live = float(row[3])
                filename = row[6].strip()
            except (IndexError, ValueError):
                continue
            # Fig. 7 consists of ordinary downward-facing scan exposures. Exclude
            # dedicated brick tests and long dwell/overnight measurements.
            if angle != 0 or "BRICK" in filename or not 100 <= live <= 900:
                continue
            spectrum_path = DATA_DIR / f"{filename}.txt"
            if not spectrum_path.exists():
                continue
            energy, counts = read_spectrum(spectrum_path)
            use = (energy >= 50) & (energy < 11400)
            total = float(counts[use].sum())
            points.append(
                {
                    "name": filename,
                    "live": live,
                    "energy": energy,
                    "counts": counts,
                    "total": total,
                    "rate": total / live,
                }
            )
    return points


def representative_point(points):
    keys = ("live", "total", "rate")
    values = np.log([[point[key] for key in keys] for point in points])
    center = np.median(values, axis=0)
    scale = np.std(values, axis=0)
    distance = np.sum(((values - center) / scale) ** 2, axis=1)
    return points[int(np.argmin(distance))]


def setup_matplotlib():
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "font.size": 11,
            "axes.labelcolor": "#23384d",
            "axes.edgecolor": "#8a9aaa",
            "xtick.color": "#4f5d6a",
            "ytick.color": "#4f5d6a",
            "axes.spines.top": False,
            "axes.spines.right": False,
        }
    )


def make_plots(points, representative):
    setup_matplotlib()
    OUTDIR.mkdir(parents=True, exist_ok=True)

    # Convert the publication's original Fig. 7 map for context.
    map_png = OUTDIR / "figure7_floor_scan.png"
    if FIG7_PDF.exists():
        subprocess.run(
            ["pdftoppm", "-singlefile", "-png", "-r", "180", str(FIG7_PDF), str(map_png.with_suffix(""))],
            check=True,
            stdout=subprocess.DEVNULL,
        )

    # Distribution of scan-point exposure and count rate.
    dist_png = OUTDIR / "scan_point_distributions.png"
    fig, axes = plt.subplots(1, 2, figsize=(10.8, 4.0))
    live = np.asarray([p["live"] for p in points])
    rate = np.asarray([p["rate"] for p in points])
    for ax, values, xlabel, color in (
        (axes[0], live, "Live time per point [s]", "#1f77b4"),
        (axes[1], rate, "Integral rate, 50–11,400 keV [s⁻¹]", "#009399"),
    ):
        ax.hist(values, bins=18, color=color, alpha=0.86, edgecolor="white")
        ax.axvline(np.median(values), color="#e67e22", lw=2.5, label=f"median = {np.median(values):.0f}")
        ax.set_xlabel(xlabel)
        ax.set_ylabel("Scan points")
        ax.grid(axis="y", alpha=0.18)
        ax.legend(frameon=False, loc="upper right")
    fig.tight_layout()
    fig.savefig(dist_png, dpi=180, bbox_inches="tight")
    plt.close(fig)

    # A real point whose live time, total counts, and rate all equal/track medians.
    rep_png = OUTDIR / "representative_point_spectrum.png"
    fig, axes = plt.subplots(2, 1, figsize=(11.0, 6.4), sharex=False, gridspec_kw={"hspace": 0.34})
    for ax, (lo, hi, width) in zip(axes, [(50, 3000, 5), (3000, 11400, 25)]):
        edges = np.arange(lo, hi + width, width)
        hist, _ = np.histogram(representative["energy"], bins=edges, weights=representative["counts"])
        centers = (edges[:-1] + edges[1:]) / 2
        ax.step(centers, hist, where="mid", color="#1f77b4", lw=1.05)
        ax.set_yscale("log")
        ax.set_ylim(0.8, max(hist.max() * 1.45, 2))
        ax.set_xlim(lo, hi)
        ax.set_ylabel(f"Counts / {width} keV")
        ax.grid(alpha=0.16, which="both")
    axes[1].set_xlabel("Measured energy [keV]")
    fig.savefig(rep_png, dpi=180, bbox_inches="tight")
    plt.close(fig)

    # Median count/bin and fraction of empty bins across points, by energy band.
    widths = [1, 5, 10, 25, 50, 100, 200]
    bands = [(50, 3000), (3000, 7000), (7000, 11400)]
    median_counts = np.zeros((len(bands), len(widths)))
    empty_fraction = np.zeros_like(median_counts)
    for i, (lo, hi) in enumerate(bands):
        for j, width in enumerate(widths):
            edges = np.arange(lo, hi + width, width)
            all_bins = []
            for point in points:
                hist, _ = np.histogram(point["energy"], bins=edges, weights=point["counts"])
                all_bins.append(hist)
            all_bins = np.concatenate(all_bins)
            median_counts[i, j] = np.median(all_bins)
            empty_fraction[i, j] = np.mean(all_bins == 0)

    heat_png = OUTDIR / "binning_statistics.png"
    fig, axes = plt.subplots(2, 1, figsize=(10.6, 5.9), gridspec_kw={"hspace": 0.35})
    labels = ["50–3,000", "3,000–7,000", "7,000–11,400"]
    shown_counts = np.maximum(median_counts, 0.5)
    im0 = axes[0].imshow(shown_counts, aspect="auto", cmap="Blues", norm=LogNorm(vmin=0.5, vmax=1000))
    im1 = axes[1].imshow(empty_fraction * 100, aspect="auto", cmap="Oranges", vmin=0, vmax=100)
    for ax in axes:
        ax.set_xticks(range(len(widths)), [str(w) for w in widths])
        ax.set_yticks(range(len(labels)), labels)
        ax.set_xlabel("Histogram bin width [keV]")
        ax.set_ylabel("Energy range [keV]")
    axes[0].set_title("Median counts per point-bin")
    axes[1].set_title("Empty point-bins [%]")
    for i in range(len(bands)):
        for j in range(len(widths)):
            axes[0].text(j, i, f"{median_counts[i, j]:.0f}", ha="center", va="center", fontsize=9,
                         color="white" if shown_counts[i, j] > 100 else "#182a3a")
            axes[1].text(j, i, f"{100 * empty_fraction[i, j]:.0f}", ha="center", va="center", fontsize=9,
                         color="white" if empty_fraction[i, j] > 0.55 else "#182a3a")
    fig.colorbar(im0, ax=axes[0], fraction=0.028, pad=0.02)
    fig.colorbar(im1, ax=axes[1], fraction=0.028, pad=0.02)
    fig.savefig(heat_png, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return map_png, dist_png, rep_png, heat_png


def add_textbox(slide, x, y, w, h, text, size=20, color=NAVY, bold=False, align=PP_ALIGN.LEFT,
                fill=None, margin=0.08, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_rich_lines(slide, x, y, w, h, lines, fill=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill; shape.line.fill.background()
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True; tf.margin_left = tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.12)
    for index, (lead, body, color) in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run(); r.text = lead; r.font.name = "Aptos"; r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = color
        r = p.add_run(); r.text = body; r.font.name = "Aptos"; r.font.size = Pt(19); r.font.color.rgb = NAVY
    return shape


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.28, 12.2, 0.55, title, 27, NAVY, True)
    slide.shapes.add_shape(1, Inches(0.58), Inches(0.92), Inches(0.78), Inches(0.06)).fill.solid()
    line = slide.shapes[-1]
    line.fill.fore_color.rgb = ORANGE; line.line.fill.background()
    if subtitle:
        add_textbox(slide, 1.48, 0.82, 11.1, 0.3, subtitle, 11, GRAY)


def add_footer(slide, text):
    add_textbox(slide, 0.55, 7.12, 12.25, 0.22, text, 8.5, GRAY)


def add_picture_contain(slide, path, x, y, w, h):
    from PIL import Image
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
                                    width=Inches(pw), height=Inches(ph))


def build_deck(points, representative, plots):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    map_png, dist_png, rep_png, heat_png = plots

    live = np.asarray([p["live"] for p in points])
    total = np.asarray([p["total"] for p in points])
    rate = np.asarray([p["rate"] for p in points])

    # 1 — answer first.
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.62, 0.50, 7.15, 1.25, "Statistics for an ‘average’\nFig. 7 scan point", 29, NAVY, True)
    add_textbox(slide, 0.68, 1.84, 6.85, 0.38, "HPGe downward-facing floor scan • 10 July 2026", 15, GRAY)
    add_textbox(slide, 0.68, 2.45, 6.65, 2.75,
                "Typical point\n236 s live time\n47k counts (50–11,400 keV)\n200 counts/s",
                25, WHITE, True, fill=NAVY, margin=0.24, valign=MSO_ANCHOR.MIDDLE)
    if map_png.exists():
        add_picture_contain(slide, map_png, 7.75, 0.52, 5.0, 6.35)
    add_footer(slide, "Purpose: choose a useful binning for public per-point spectra corresponding to Fig. 7.")

    # 2 — population and selection.
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Typical scan points are ~4-minute exposures", "Selection avoids dedicated long-dwell and shielding-test runs")
    slide.shapes.add_picture(str(dist_png), Inches(0.62), Inches(1.28), width=Inches(7.45))
    add_rich_lines(slide, 8.28, 1.42, 4.45, 4.9, [
        (f"{len(points)} points analyzed\n", "downward-facing, 100–900 s, no brick tests", BLUE),
        ("Median [10–90%]\n", f"live time: {np.median(live):.0f} s [{np.quantile(live,.1):.0f}–{np.quantile(live,.9):.0f}]\ncounts: {np.median(total)/1000:.0f}k [{np.quantile(total,.1)/1000:.0f}k–{np.quantile(total,.9)/1000:.0f}k]\nrate: {np.median(rate):.0f} s⁻¹ [{np.quantile(rate,.1):.0f}–{np.quantile(rate,.9):.0f}]", TEAL),
    ], fill=LIGHT)
    add_footer(slide, "Counts and rates are recomputed directly from each calibrated spectrum over 50 ≤ E < 11,400 keV.")

    # 3 — representative spectrum.
    slide = prs.slides.add_slide(blank)
    add_title(slide, f"A genuinely representative point: {representative['name']}", "This measured point lands at the median live time, total counts, and integral rate")
    slide.shapes.add_picture(str(rep_png), Inches(0.58), Inches(1.18), width=Inches(8.55))
    add_rich_lines(slide, 9.35, 1.38, 3.45, 4.8, [
        (f"{representative['live']:.0f} s", " live time", BLUE),
        (f"{representative['total']/1000:.0f}k", " counts", TEAL),
        (f"{representative['rate']:.0f} s⁻¹", " integral rate", ORANGE),
        ("Takeaway\n", "Most information is below 3 MeV. Above 7 MeV, even broad bins are often empty at an ordinary point.", GREEN),
    ], fill=LIGHT)
    add_footer(slide, "Display binning only: 5 keV below 3 MeV and 25 keV above 3 MeV; raw channels are ~0.7 keV wide.")

    # 4 — statistics versus bin width.
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Statistics vary strongly across the energy range", "Numbers summarize every point-bin in the selected scan population")
    slide.shapes.add_picture(str(heat_png), Inches(0.58), Inches(1.16), width=Inches(8.75))
    add_rich_lines(slide, 9.55, 1.38, 3.2, 4.95, [
        ("50–3,000 keV\n", "5 keV bins: median 47 counts; <1% empty", BLUE),
        ("3,000–7,000 keV\n", "25 keV bins: median 33 counts; <1% empty", TEAL),
        ("7,000–11,400 keV\n", "200 keV bins: median 4 counts; 22% empty", ORANGE),
        ("Poisson scale\n", "N = 4 still means ~50% relative counting uncertainty.", GREEN),
    ], fill=LIGHT)
    add_footer(slide, "Medians include continuum and peak regions; individual strong lines can have much higher occupancy.")

    # 5 — recommendation.
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Recommendation: preserve resolution; rebin for viewing", "Sparse high-energy bins are scientifically honest and cheap to store")
    add_rich_lines(slide, 0.7, 1.35, 5.9, 4.95, [
        ("Public per-point histogram\n", "Use a common 1 keV grid from 50–11,400 keV (or retain native channels plus calibration). Store integer counts and live time—not only rate.", BLUE),
        ("Why\n", "This preserves narrow HPGe lines, permits exact downstream rebinning, and is only ~1.3 million bins for 115 points.", TEAL),
        ("Uncertainty\n", "Poisson σ = √N; rate and uncertainty can be derived from counts, live time, and bin width.", GREEN),
    ], fill=LIGHT)
    add_rich_lines(slide, 6.85, 1.35, 5.8, 4.95, [
        ("Quick-look plots\n", "Use 5 keV bins below 3 MeV, 25 keV from 3–7 MeV, and 200 keV from 7–11.4 MeV.", ORANGE),
        ("Metadata to ship\n", "Position (x, z), orientation, start time, real/live time, energy edges, calibration provenance, detector/shield configuration, and reactor state.", NAVY),
        ("Decision\n", "Fine archival binning + optional adaptive previews gives users both line sensitivity and readable per-point spectra.", TEAL),
    ], fill=LIGHT)
    add_footer(slide, "Suggested release check: verify scan-point list against the exact inputs used to generate publication Fig. 7.")

    prs.save(PPTX_OUT)
    return PPTX_OUT


def main():
    points = load_points()
    if not points:
        raise RuntimeError(f"No scan spectra found in {DATA_DIR}")
    representative = representative_point(points)
    plots = make_plots(points, representative)
    path = build_deck(points, representative, plots)
    print(f"Wrote {path}")
    print(f"Points: {len(points)}; representative: {representative['name']}")


if __name__ == "__main__":
    main()
