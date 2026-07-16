#!/usr/bin/env python3
from pathlib import Path
import csv

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

REPO = Path('/home/blaine/projects/HFIR_BG_Analysis')
COMPARE = REPO / 'analysis/unfold/algorithm_comparison'
OUTDIR = REPO / 'analysis/unfold/slides'
OUTDIR.mkdir(parents=True, exist_ok=True)
OUTFILE = OUTDIR / 'unfolding_algorithm_comparison_poisson_pgd_2026-04-16.pptx'

ASSETDIR = OUTDIR / 'formula_assets'
ASSETDIR.mkdir(parents=True, exist_ok=True)

BG = RGBColor(248, 246, 241)
ACCENT = RGBColor(120, 44, 44)
TEXT = RGBColor(34, 34, 34)
MUTED = RGBColor(90, 90, 90)
TABLE_HEAD = RGBColor(232, 224, 219)
TABLE_BODY = RGBColor(252, 251, 248)


def style_slide(slide, slide_number, total_slides):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.color.rgb = BG
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12.2), Inches(0.03))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.color.rgb = ACCENT

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(11.4), Inches(0.22))
    p = footer.text_frame.paragraphs[0]
    p.text = 'HFIR gamma paper | unfolding algorithm crosscheck'
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED

    num = slide.shapes.add_textbox(Inches(12.1), Inches(6.95), Inches(0.7), Inches(0.25))
    p2 = num.text_frame.paragraphs[0]
    p2.text = f'{slide_number}/{total_slides}'
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(10)
    p2.font.color.rgb = MUTED


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.2), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TEXT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.8), Inches(12.0), Inches(0.35))
        p2 = sub_box.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED


def add_bullets(slide, bullets, left, top, width, height, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(6)


def add_image(slide, path, left, top, width=None, height=None):
    kwargs = {}
    if width is not None:
        kwargs['width'] = Inches(width)
    if height is not None:
        kwargs['height'] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)




def render_formula(formula, outfile, fontsize=24, dpi=300, pad=0.2):
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.patch.set_alpha(0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis('off')
    text_artist = ax.text(0, 0.5, formula, fontsize=fontsize, color='#222222', va='center', ha='left')
    fig.canvas.draw()
    bbox = text_artist.get_window_extent(renderer=fig.canvas.get_renderer()).expanded(1.04, 1.18)
    width = max(0.01, bbox.width / dpi + pad)
    height = max(0.01, bbox.height / dpi + pad)
    plt.close(fig)

    fig = plt.figure(figsize=(width, height))
    fig.patch.set_alpha(0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis('off')
    ax.text(0.01, 0.5, formula, fontsize=fontsize, color='#222222', va='center', ha='left')
    fig.savefig(outfile, dpi=dpi, transparent=True, bbox_inches='tight', pad_inches=0.04)
    plt.close(fig)


def add_formula(slide, formula, left, top, width, fontsize=24, name='formula'):
    outfile = ASSETDIR / f'{name}.png'
    render_formula(formula, outfile, fontsize=fontsize)
    slide.shapes.add_picture(str(outfile), Inches(left), Inches(top), width=Inches(width))

def add_table(slide, data, left, top, width, height, font_size=15):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size if r else font_size + 1)
                p.font.color.rgb = TEXT
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                if r == 0:
                    p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEAD if r == 0 else TABLE_BODY
    return table


def load_rows():
    path = COMPARE / 'algorithm_pairwise_comparison.csv'
    with open(path) as f:
        rows = list(csv.DictReader(f))
    return rows


def find_row(rows, response, filename):
    for row in rows:
        if row['response'] == response and row['filename'] == filename:
            return row
    raise KeyError((response, filename))


def fmt(x, nd=3):
    return f'{float(x):.{nd}f}'


def build_deck():
    rows = load_rows()
    mif_iso = find_row(rows, 'isotropic', 'MIF_BOX_REACTOR_OPTIMIZED_DAYCOUNT_OPTIMAL_GAIN')
    hb4_iso = find_row(rows, 'isotropic', 'HB4_DOWN_OVERNIGHT_1')
    mif_front = find_row(rows, 'front', 'MIF_BOX_REACTOR_OPTIMIZED_DAYCOUNT_OPTIMAL_GAIN')
    hb4_front = find_row(rows, 'front', 'HB4_DOWN_OVERNIGHT_1')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total_slides = 6

    slide = prs.slides.add_slide(blank)
    style_slide(slide, 1, total_slides)
    add_title(slide, 'Unfolding algorithm crosscheck', 'Same six spectra, same migration matrices, RL vs penalized Poisson mirror descent')
    add_bullets(
        slide,
        [
            'Question: is the unrealistically low 0.05-2 MeV unfolded flux at reactor-facing MIF an artifact of Richardson-Lucy, or does it come from the response model?',
            'Crosscheck performed: existing Richardson-Lucy (RL) vs a new non-negative penalized Poisson mirror-descent solver.',
            'Both algorithms were run on the same six paper spectra with both isotropic and front-face migration matrices.',
            'Takeaway preview: the low-energy suppression is not specific to RL. It persists across solvers, which points to migration-matrix / detector-response mismodeling rather than to the inversion method alone.'
        ],
        left=0.7, top=1.35, width=6.0, height=4.8, font_size=19,
    )
    add_table(
        slide,
        [
            ['Case', 'RL low frac', 'Alt low frac'],
            ['MIF isotropic', fmt(mif_iso['low_fraction_rl'], 4), fmt(mif_iso['low_fraction_alt'], 4)],
            ['MIF front', fmt(mif_front['low_fraction_rl'], 4), fmt(mif_front['low_fraction_alt'], 4)],
            ['HB4 isotropic', fmt(hb4_iso['low_fraction_rl'], 4), fmt(hb4_iso['low_fraction_alt'], 4)],
            ['HB4 front', fmt(hb4_front['low_fraction_rl'], 4), fmt(hb4_front['low_fraction_alt'], 4)],
        ],
        left=7.1, top=1.55, width=5.5, height=2.6, font_size=15,
    )
    add_image(slide, COMPARE / 'low_energy_fraction_comparison.png', left=7.0, top=4.3, width=5.7)

    slide = prs.slides.add_slide(blank)
    style_slide(slide, 2, total_slides)
    add_title(slide, 'Methods compared', 'Definitions and update equations for the two unfold solvers')
    add_bullets(
        slide,
        [
            'State variables',
            'A_ij: migration-matrix probability for true-energy bin j to contribute to measured-energy bin i.',
            'y_i: measured counts in detector-energy bin i. x_j: unfolded incident flux in true-energy bin j.',
            '(Ax)_i: refolded prediction in measured bin i. n: iteration index.',
            'Solver controls',
            'lambda: smoothness-penalty strength in the Poisson-PGD objective. eta: mirror-descent step size.',
            'Interpretation',
            'RL uses only the Poisson data term and is regularized mainly by stopping early. Poisson-PGD uses the same data term plus an explicit smoothness penalty.',
            'Method background: mirror descent from Beck and Teboulle (2003); penalized Poisson inverse problems from Bardsley and Luttman (2009).'
        ],
        left=0.7, top=1.18, width=5.3, height=5.7, font_size=17,
    )
    add_formula(slide, r'$x_j^{(n+1)} = x_j^{(n)} \frac{\sum_i A_{ij} \, y_i / (A x^{(n)})_i}{\sum_i A_{ij}}$', left=6.1, top=1.55, width=6.0, fontsize=28, name='rl_update')
    add_formula(slide, r'$\Phi(x) = \sum_i \left[(Ax)_i - y_i \ln (Ax)_i\right] + \lambda \sum_j (x_j - x_{j-1})^2$', left=5.95, top=2.95, width=6.15, fontsize=28, name='pgd_objective')
    add_formula(slide, r'$g_j = \sum_i A_{ij}\left(1 - y_i/(Ax^{(n)})_i\right) + 2\lambda \left(2x_j - x_{j-1} - x_{j+1}\right)$', left=5.95, top=4.15, width=6.2, fontsize=28, name='pgd_gradient')
    add_formula(slide, r'$x_j^{(n+1)} = x_j^{(n)} \exp\!\left[-\eta \, g_j / \sum_i A_{ij}\right]$', left=6.15, top=5.45, width=5.95, fontsize=28, name='pgd_update')

    slide = prs.slides.add_slide(blank)
    style_slide(slide, 3, total_slides)
    add_title(slide, 'Isotropic matrix: dip remains with a statistically matched alternate solver')
    add_bullets(
        slide,
        [
            f"MIF reactor-on isotropic: RL low fraction = {fmt(mif_iso['low_fraction_rl'], 4)}, Poisson-PGD = {fmt(mif_iso['low_fraction_alt'], 4)}. Both are suppressed.",
            f"HB4 isotropic: RL low fraction = {fmt(hb4_iso['low_fraction_rl'], 4)}, Poisson-PGD = {fmt(hb4_iso['low_fraction_alt'], 4)}. Same qualitative trend.",
            f"For MIF isotropic, Alt/RL low-flux ratio = {fmt(mif_iso['low_flux_ratio_alt_over_rl'], 3)}.",
            'Interpretation: replacing RL does not remove the dip, so the dominant problem is in the migration matrix and the detector-shield response model, not in the inverse solver by itself.'
        ],
        left=0.7, top=1.25, width=5.5, height=5.3, font_size=18,
    )
    add_image(slide, COMPARE / 'algorithm_overlay_isotropic.png', left=6.35, top=1.25, width=6.35)

    slide = prs.slides.add_slide(blank)
    style_slide(slide, 4, total_slides)
    add_title(slide, 'Front matrix: solution is very stable across solvers')
    add_bullets(
        slide,
        [
            f"MIF reactor-on front: RL low fraction = {fmt(mif_front['low_fraction_rl'], 4)}, Poisson-PGD = {fmt(mif_front['low_fraction_alt'], 4)}.",
            f"HB4 front: RL low fraction = {fmt(hb4_front['low_fraction_rl'], 4)}, Poisson-PGD = {fmt(hb4_front['low_fraction_alt'], 4)}.",
            f"For MIF front, Alt/RL low-flux ratio = {fmt(mif_front['low_flux_ratio_alt_over_rl'], 3)}; for HB4 front it is {fmt(hb4_front['low_flux_ratio_alt_over_rl'], 3)}.",
            'Interpretation: the front-face matrix reduces the effect, but does not remove all low-energy suppression. That points to remaining detector, collimator, or shielding mismodeling in the response simulation.'
        ],
        left=0.7, top=1.25, width=5.5, height=5.2, font_size=18,
    )
    add_image(slide, COMPARE / 'algorithm_overlay_front.png', left=6.35, top=1.25, width=6.35)

    slide = prs.slides.add_slide(blank)
    style_slide(slide, 5, total_slides)
    add_title(slide, 'Toy-MC spread', 'Point-estimate comparison plus toy-counting fluctuations on the key directional cases')
    add_bullets(
        slide,
        [
            'We do not yet have full unfolded covariance matrices. Instead, I fluctuated the original GeDataHist counts and reran both unfold algorithms on matched toys.',
            'MIF reactor-on: solver shift at isotropic is -0.0032 ± 0.0014 in low-energy fraction; solver shift at front is +0.0004 ± 0.0002.',
            'MIF reactor-on: front minus isotropic is +0.0467 ± 0.0029 for RL and +0.0503 ± 0.0038 for Poisson-PGD.',
            'HB4: solver shift at isotropic is -0.0414 ± 0.0010; solver shift at front is +0.0004 ± 0.0000.',
            'HB4: front minus isotropic is +0.0300 ± 0.0005 for RL and +0.0718 ± 0.0007 for Poisson-PGD.'
        ],
        left=0.7, top=1.35, width=6.1, height=5.1, font_size=18,
    )
    add_table(
        slide,
        [
            ['Case', 'Algorithm diff', 'Response diff'],
            ['MIF Rx On', '-0.0032 ± 0.0014 (iso)\n+0.0004 ± 0.0002 (front)', '+0.0467 ± 0.0029 (RL)\n+0.0503 ± 0.0038 (PGD)'],
            ['HB4', '-0.0414 ± 0.0010 (iso)\n+0.0004 ± 0.0000 (front)', '+0.0300 ± 0.0005 (RL)\n+0.0718 ± 0.0007 (PGD)'],
        ],
        left=7.05, top=1.8, width=5.5, height=2.6, font_size=14,
    )
    add_bullets(
        slide,
        [
            'Reading: for reactor-facing MIF, the front-vs-isotropic shift is larger than the solver shift, but the residual front-case suppression shows that geometry alone is not the full story. The migration-matrix model remains the limiting issue.'
        ],
        left=7.05, top=4.7, width=5.5, height=1.3, font_size=15,
    )

    slide = prs.slides.add_slide(blank)
    style_slide(slide, 6, total_slides)
    add_title(slide, 'Conclusion for the paper', 'What this crosscheck now supports and what it still does not')
    add_bullets(
        slide,
        [
            'The low-energy suppression is not an RL-only artifact. It persists with a second solver built around the same Poisson counting model.',
            'The strongest evidence points to migration-matrix mismatch in the detector-response simulation, including detector, collimator, shielding, and illumination modeling.',
            'The front-face matrix improves the situation and is solver-stable, but some low-energy suppression remains, so the problem is not solved by angular reweighting alone.',
            'We still do not have a full unfolded covariance matrix, but the key MIF and HB4 statements now include toy-MC counting-statistics spreads.',
            'Recommendation: describe the alternate Poisson-PGD crosscheck as evidence that the dominant uncertainty lies in the migration-matrix / detector-response simulation, not in Richardson-Lucy itself.'
        ],
        left=0.7, top=1.35, width=6.0, height=5.1, font_size=19,
    )
    add_table(
        slide,
        [
            ['Metric', 'MIF isotropic', 'MIF front'],
            ['RL low fraction', fmt(mif_iso['low_fraction_rl'], 4), fmt(mif_front['low_fraction_rl'], 4)],
            ['PGD low fraction', fmt(mif_iso['low_fraction_alt'], 4), fmt(mif_front['low_fraction_alt'], 4)],
            ['RL refold chi2', fmt(mif_iso['refold_chi2_rl'], 1), fmt(mif_front['refold_chi2_rl'], 1)],
            ['Alt refold chi2', fmt(mif_iso['refold_chi2_alt'], 1), fmt(mif_front['refold_chi2_alt'], 1)],
        ],
        left=7.05, top=1.65, width=5.6, height=2.8, font_size=15,
    )
    add_image(slide, COMPARE / 'low_energy_fraction_comparison.png', left=7.0, top=4.5, width=5.7)

    prs.save(OUTFILE)
    print(f'Wrote {OUTFILE}')


if __name__ == '__main__':
    build_deck()
