#!/usr/bin/env python3
import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


REPO = Path(__file__).resolve().parents[1]
ANALYSIS_ROOT = Path(os.environ.get("HFIRBG_ANALYSIS", REPO / "analysis")).expanduser()
PAPER_FILES = ANALYSIS_ROOT / "unfold/paper_files"
OUTDIR = ANALYSIS_ROOT / "unfold/slides"
OUTDIR.mkdir(parents=True, exist_ok=True)


BG = RGBColor(248, 246, 241)
ACCENT = RGBColor(120, 44, 44)
TEXT = RGBColor(34, 34, 34)
MUTED = RGBColor(90, 90, 90)


def style_slide(slide, slide_number, total_slides):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.color.rgb = BG
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12.2), Inches(0.03))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.color.rgb = ACCENT

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(11.4), Inches(0.22))
    p = footer.text_frame.paragraphs[0]
    p.text = "HFIR gamma paper update | front-face response status"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED

    num = slide.shapes.add_textbox(Inches(12.1), Inches(6.95), Inches(0.7), Inches(0.25))
    p2 = num.text_frame.paragraphs[0]
    p2.text = f"{slide_number}/{total_slides}"
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(10)
    p2.font.color.rgb = MUTED


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TEXT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.8), Inches(12.0), Inches(0.35))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED


def add_bullets(slide, bullets, left=0.6, top=1.2, width=5.9, height=5.6, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.level = 0
        p.space_after = Pt(8)


def add_image(slide, path, left, top, width=None, height=None):
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def add_table(slide, rows, cols, data, left, top, width, height):
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(16 if r == 0 else 15)
                p.font.color.rgb = TEXT
                if r == 0:
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(232, 224, 219)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(252, 251, 248)
    return table


def make_needed_stats_plot(csv_path, out_path):
    df = pd.read_csv(csv_path)
    df["additional_stats_factor_needed"] = df["front_to_iso_median_error_ratio"] ** 2

    fig, ax = plt.subplots(figsize=(9, 4.8))
    ax.axhline(1.0, color="black", linestyle="--", linewidth=1)
    ax.plot(df["energy_keV"], df["additional_stats_factor_needed"], marker="o", linewidth=1.4, color="#d62728")
    ax.set_yscale("log")
    ax.set_xlabel("Energy [keV]")
    ax.set_ylabel("Extra front stats factor\nneeded to match isotropic")
    ax.set_title("Front production still needed at energies already simulated")
    ax.grid(True, alpha=0.3)
    fig.tight_layout()
    fig.savefig(out_path, dpi=200, bbox_inches="tight")
    fig.savefig(out_path.with_suffix(".pdf"), bbox_inches="tight")
    plt.close(fig)
    return df


def make_runtime_plot(csv_path, out_path):
    df = pd.read_csv(csv_path)
    fig, ax = plt.subplots(figsize=(9, 4.8))
    ax.plot(df["energy_keV"], df["iso_runtime_s"] / 3600.0, marker="o", linewidth=1.2, label="Isotropic")
    ax.plot(df["energy_keV"], df["front_runtime_s"] / 3600.0, marker="o", linewidth=1.2, label="Front")
    ax.set_xlabel("Energy [keV]")
    ax.set_ylabel("Combined runtime [hours]")
    ax.set_title("Current simulated exposure at common energies")
    ax.set_yscale("log")
    ax.grid(True, alpha=0.3)
    ax.legend()
    fig.tight_layout()
    fig.savefig(out_path, dpi=200, bbox_inches="tight")
    fig.savefig(out_path.with_suffix(".pdf"), bbox_inches="tight")
    plt.close(fig)


def build_deck():
    stats_csv = PAPER_FILES / "migration_matrix_stats_comparison.csv"
    support_txt = PAPER_FILES / "migration_matrix_support_comparison.txt"
    support_plot = PAPER_FILES / "migration_matrix_support_comparison.png"
    bounds_plot = PAPER_FILES / "unfolded_spectrum_bounds.png"
    iso_plot = PAPER_FILES / "all_unfolded_spectra_isotropic.png"
    front_plot = PAPER_FILES / "all_unfolded_spectra_front.png"

    needed_plot = OUTDIR / "front_stats_needed_factor.png"
    runtime_plot = OUTDIR / "front_vs_iso_runtime.png"
    df = make_needed_stats_plot(stats_csv, needed_plot)
    make_runtime_plot(stats_csv, runtime_plot)

    support_summary = support_txt.read_text().strip().splitlines()
    n_common = len(df)
    n_missing = 166 - 17
    worst_needed = df["additional_stats_factor_needed"].max()
    low_energy_needed = df[df["energy_keV"] <= 10400]["additional_stats_factor_needed"].median()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    total_slides = 7
    slide_index = 1

    slide = prs.slides.add_slide(blank)
    style_slide(slide, slide_index, total_slides)
    add_title(slide, "Front unfolding status update", "HFIR gamma paper response-model crosscheck, March 12 2026")
    add_bullets(
        slide,
        [
            "This is a status check, not a final front-sim result.",
            "I rebuilt the front-face migration matrix from ~/Downloads/2026-03-11.zip, reran the six unfolds, and compared them to the existing isotropic case.",
            "The two questions for today are simple: what do the partial front sims already show, and how much more do we still need before swapping them into the paper with confidence?",
        ],
        left=0.7,
        top=1.35,
        width=6.1,
        height=4.8,
        font_size=20,
    )
    add_image(slide, support_plot, left=7.1, top=1.2, width=5.7)

    slide_index += 1
    slide = prs.slides.add_slide(blank)
    style_slide(slide, slide_index, total_slides)
    add_title(slide, "Numbers to quote")
    quote_table = [
        ["Quantity", "Current value"],
        ["Isotropic direct-support energies", "166"],
        ["Front direct-support energies", "17"],
        ["Missing front energies vs isotropic", f"{n_missing}"],
        ["Largest front gap to nearest direct sim", "4.46 MeV"],
        ["Typical extra front stats still needed", f"{low_energy_needed:.1f}x"],
        ["Worst shared-energy stats deficit", f"{worst_needed:.1f}x"],
    ]
    add_table(slide, len(quote_table), 2, quote_table, left=0.7, top=1.35, width=5.3, height=4.7)
    add_bullets(
        slide,
        [
            "Bottom line: enough to show the direction of the effect, not enough to call the front matrix finished.",
            "The weak spot is the low and middle of the spectrum. The high-energy cluster is in much better shape.",
        ],
        left=0.8,
        top=6.0,
        width=5.3,
        height=0.7,
        font_size=16,
    )
    add_image(slide, support_plot, left=6.35, top=1.15, width=6.3)

    slide_index += 1
    slide = prs.slides.add_slide(blank)
    style_slide(slide, slide_index, total_slides)
    add_title(slide, "How much more front statistics are still needed?")
    add_bullets(
        slide,
        [
            "This is the cleanest collaborator-facing estimate I could make from the shared energies.",
            "For each common energy, I asked how much more front exposure would be needed to match the isotropic median bin error.",
            f"Up to about 10.4 MeV, the front production is typically short by roughly {low_energy_needed:.1f}x.",
            f"The worst common-energy deficit is about {worst_needed:.1f}x.",
            "By 10.5 to 10.8 MeV, we are already in much better shape. The unfinished part is mostly the low and middle of the spectrum.",
        ],
        left=0.7,
        top=1.2,
        width=5.5,
        height=5.4,
        font_size=17,
    )
    add_image(slide, needed_plot, left=6.35, top=1.2, width=6.3)

    slide_index += 1
    slide = prs.slides.add_slide(blank)
    style_slide(slide, slide_index, total_slides)
    add_title(slide, "Current exposure at energies both productions share")
    add_bullets(
        slide,
        [
            "The front campaign is uneven. Around 1 MeV and 10.0 to 10.4 MeV it is still thin.",
            "Once the big 10.5 to 10.8 MeV production starts, the front case catches up fast and in some bins passes the isotropic exposure.",
            "So the high-energy end of the bounds is already on firmer ground than the low-energy end.",
        ],
        left=0.7,
        top=1.2,
        width=5.5,
        height=4.5,
        font_size=19,
    )
    add_image(slide, runtime_plot, left=6.35, top=1.2, width=6.3)

    slide_index += 1
    slide = prs.slides.add_slide(blank)
    style_slide(slide, slide_index, total_slides)
    add_title(slide, "Partial physics result: isotropic vs front bounds")
    add_bullets(
        slide,
        [
            "The front case is not just adding noise. It changes the unfold in a consistent way.",
            "At five of the six locations, it pulls the 2 to 9 MeV continuum below the isotropic result.",
            "At the reactor-facing MIF point, the real story is lower in energy: the front response removes the weird 0 to 2 MeV dip in the isotropic unfold.",
            "I read that as a geometry mismatch in the isotropic approximation for the most directional measurement, not as a general problem with the detector or collimator model.",
        ],
        left=0.7,
        top=1.15,
        width=5.2,
        height=5.8,
        font_size=18,
    )
    add_image(slide, bounds_plot, left=5.95, top=1.1, width=6.8)

    slide_index += 1
    slide = prs.slides.add_slide(blank)
    style_slide(slide, slide_index, total_slides)
    add_title(slide, "Overlay of current unfolded spectra")
    add_image(slide, iso_plot, left=0.55, top=1.2, width=6.1)
    add_image(slide, front_plot, left=6.7, top=1.2, width=6.1)
    cap1 = slide.shapes.add_textbox(Inches(0.75), Inches(6.7), Inches(5.5), Inches(0.3))
    cap1.text_frame.text = "Isotropic response"
    cap1.text_frame.paragraphs[0].font.size = Pt(16)
    cap1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cap1.text_frame.paragraphs[0].font.color.rgb = TEXT
    cap2 = slide.shapes.add_textbox(Inches(6.9), Inches(6.7), Inches(5.5), Inches(0.3))
    cap2.text_frame.text = "Front-face response"
    cap2.text_frame.paragraphs[0].font.size = Pt(16)
    cap2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cap2.text_frame.paragraphs[0].font.color.rgb = TEXT

    slide_index += 1
    slide = prs.slides.add_slide(blank)
    style_slide(slide, slide_index, total_slides)
    add_title(slide, "Recommendation for the paper")
    add_bullets(
        slide,
        [
            "I think the paper can stay basically final now.",
            "Use the current front case as a directional limiting case and keep the bounds figure in the draft.",
            "Do not sell the front matrix as complete. The support plot makes that caveat easy to show.",
            "When the remaining front jobs land, we only need to rebuild the matrix, rerun the six unfolds, swap the plots, and revisit the wording if the bounds actually move.",
        ],
        left=0.8,
        top=1.35,
        width=11.5,
        height=4.8,
        font_size=21,
    )

    out = OUTDIR / "front_unfolding_update_2026-03-12.pptx"
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    build_deck()
