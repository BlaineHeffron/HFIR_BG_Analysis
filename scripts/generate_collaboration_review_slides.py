#!/usr/bin/env python3
"""Generate collaboration slides and a printable student quickstart handout."""

from __future__ import annotations

import argparse
from functools import lru_cache
import json
import textwrap
from pathlib import Path

from matplotlib import font_manager
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = REPO_ROOT / "analysis" / "presentations" / "bryce_collaboration_review"
WIDTH, HEIGHT = 1600, 900

NAVY = "#14324a"
BLUE = "#247ba0"
GREEN = "#4f8a5b"
GOLD = "#d9a441"
RED = "#a94442"
TEXT = "#202a33"
MUTED = "#64717d"
LIGHT = "#f4f7f9"
PALE_BLUE = "#eaf3f8"
PALE_GREEN = "#edf5ef"
WHITE = "#ffffff"


@lru_cache(maxsize=None)
def font(size: int, bold: bool = False, mono: bool = False) -> ImageFont.FreeTypeFont:
    """Resolve a bundled/system font without embedding workstation paths."""
    family = "DejaVu Sans Mono" if mono else "DejaVu Sans"
    weight = "bold" if bold else "normal"
    path = font_manager.findfont(font_manager.FontProperties(family=family, weight=weight))
    return ImageFont.truetype(path, size)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument(
        "--browser-screenshot",
        type=Path,
        default=REPO_ROOT / "analysis" / "presentations" / "browser_demo.png",
    )
    return parser.parse_args()


def base_slide(title: str, subtitle: str | None, number: int, total: int) -> Image.Image:
    image = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, WIDTH, 90), fill=NAVY)
    draw.text((60, 22), title, font=font(38, bold=True), fill=WHITE)
    if subtitle:
        draw.text((62, 104), subtitle, font=font(20), fill=MUTED)
    draw.rectangle((55, 855, 1545, 857), fill=BLUE)
    draw.text((60, 865), "HFIR gamma-background data | Figure 7 and public analysis", font=font(14), fill=MUTED)
    label = f"{number}/{total}"
    box = draw.textbbox((0, 0), label, font=font(14))
    draw.text((1540 - (box[2] - box[0]), 865), label, font=font(14), fill=MUTED)
    return image


def title_slide(title: str, subtitle: str, footer: str) -> Image.Image:
    image = Image.new("RGB", (WIDTH, HEIGHT), NAVY)
    draw = ImageDraw.Draw(image)
    draw.rectangle((75, 105, 105, 700), fill=GOLD)
    draw.multiline_text((150, 155), title, font=font(62, bold=True), fill=WHITE, spacing=16)
    draw.multiline_text((155, 430), subtitle, font=font(29), fill="#d9e6ef", spacing=10)
    draw.text((155, 760), footer, font=font(20), fill="#a9bdca")
    return image


def fit_image(
    canvas: Image.Image,
    path: Path,
    box: tuple[int, int, int, int],
    border: bool = True,
    crop_fraction: tuple[float, float, float, float] | None = None,
) -> None:
    source = Image.open(path).convert("RGB")
    if crop_fraction:
        left, top, right, bottom = crop_fraction
        source = source.crop((
            round(source.width * left),
            round(source.height * top),
            round(source.width * right),
            round(source.height * bottom),
        ))
    x0, y0, x1, y1 = box
    scale = min((x1 - x0) / source.width, (y1 - y0) / source.height)
    resized = source.resize((round(source.width * scale), round(source.height * scale)), Image.Resampling.LANCZOS)
    x = x0 + (x1 - x0 - resized.width) // 2
    y = y0 + (y1 - y0 - resized.height) // 2
    canvas.paste(resized, (x, y))
    if border:
        ImageDraw.Draw(canvas).rectangle((x - 1, y - 1, x + resized.width, y + resized.height), outline="#c9d2d8", width=2)


def wrapped(draw: ImageDraw.ImageDraw, text: str, xy: tuple[int, int], width_chars: int, size: int, color: str = TEXT, bold: bool = False, spacing: int = 8) -> int:
    lines = textwrap.wrap(text, width_chars, break_long_words=False, break_on_hyphens=False)
    draw.multiline_text(xy, "\n".join(lines), font=font(size, bold=bold), fill=color, spacing=spacing)
    line_height = size + spacing
    return len(lines) * line_height


def bullets(draw: ImageDraw.ImageDraw, items: list[str], box: tuple[int, int, int, int], size: int = 27, color: str = TEXT) -> None:
    x0, y0, x1, _ = box
    y = y0
    chars = max(20, int((x1 - x0) / (size * 0.58)) - 4)
    for item in items:
        draw.ellipse((x0, y + 11, x0 + 12, y + 23), fill=BLUE)
        used = wrapped(draw, item, (x0 + 28, y), chars, size, color=color, spacing=7)
        y += used + 20


def metric_card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], value: str, label: str, accent: str = BLUE) -> None:
    draw.rounded_rectangle(box, radius=18, fill=LIGHT, outline="#d7e0e5", width=2)
    x0, y0, x1, _ = box
    draw.rectangle((x0, y0, x0 + 12, box[3]), fill=accent)
    draw.text((x0 + 35, y0 + 25), value, font=font(46, bold=True), fill=NAVY)
    wrapped(draw, label, (x0 + 36, y0 + 90), max(15, int((x1 - x0) / 16)), 19, color=MUTED)


def code_box(draw: ImageDraw.ImageDraw, code: str, box: tuple[int, int, int, int], size: int = 20) -> None:
    draw.rounded_rectangle(box, radius=15, fill="#17242e", outline="#314653", width=2)
    draw.multiline_text((box[0] + 24, box[1] + 22), code, font=font(size, mono=True), fill="#e8f0f4", spacing=9)


def simple_table(draw: ImageDraw.ImageDraw, headers: list[str], rows: list[list[str]], box: tuple[int, int, int, int], widths: list[float] | None = None) -> None:
    x0, y0, x1, y1 = box
    ncols = len(headers)
    widths = widths or [1 / ncols] * ncols
    total = sum(widths)
    positions = [x0]
    for value in widths:
        positions.append(round(positions[-1] + (x1 - x0) * value / total))
    row_height = (y1 - y0) / (len(rows) + 1)
    for row_index, row in enumerate([headers] + rows):
        top = round(y0 + row_index * row_height)
        bottom = round(y0 + (row_index + 1) * row_height)
        fill = NAVY if row_index == 0 else (WHITE if row_index % 2 else LIGHT)
        for column, value in enumerate(row):
            draw.rectangle((positions[column], top, positions[column + 1], bottom), fill=fill, outline="#cad4da", width=1)
            text_color = WHITE if row_index == 0 else TEXT
            selected_font = font(19, bold=row_index == 0)
            text_box = draw.textbbox((0, 0), str(value), font=selected_font)
            text_height = text_box[3] - text_box[1]
            draw.text((positions[column] + 12, top + (bottom - top - text_height) / 2 - 2), str(value), font=selected_font, fill=text_color)


def save_slides(slides: list[Image.Image], output_dir: Path) -> tuple[Path, Path]:
    slide_dir = output_dir / "slide_images"
    slide_dir.mkdir(parents=True, exist_ok=True)
    for stale_slide in slide_dir.glob("slide_*.png"):
        stale_slide.unlink()
    paths: list[Path] = []
    for index, slide in enumerate(slides, start=1):
        path = slide_dir / f"slide_{index:02d}.png"
        slide.save(path, optimize=True)
        paths.append(path)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for path in paths:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    pptx_path = output_dir / "HFIR_public_data_analysis_overview.pptx"
    presentation.save(pptx_path)

    pdf_path = output_dir / "HFIR_public_data_analysis_overview_preview.pdf"
    slides[0].save(pdf_path, "PDF", resolution=150, save_all=True, append_images=slides[1:])
    return pptx_path, pdf_path


def write_handout(output_dir: Path) -> tuple[Path, Path]:
    markdown = """# A beginner's guide to the HFIR gamma-background data

This guide assumes you are new to this dataset. You do **not** need to know Git, Python, or the project folders to follow the presentation or explore the browser.

## 1. Start with the idea, not the files

A **spectrum** is a record of how many gamma rays a detector observed at each energy. It is the radiation fingerprint of one place and one time.

This collection contains 1,802 calibrated spectra from 354 measurement runs at 241 mapped locations. Each spectrum is linked to basic context—where, when, and under what conditions it was measured—so comparisons can be made fairly.

## 2. Explore in the browser first

If a presenter or colleague has started the browser, open the address they provide (on the presenter's computer it is usually <http://localhost:8501>). Then:

1. Filter by a location, reactor condition, shielding setup, or a word in the run description.
2. Choose a measurement point, then select one recording from that point.
3. View the spectrum, adjust the energy range if helpful, and download the displayed values only if you want to work with them further.

You do not need to memorize run IDs or understand the repository layout to do this. A good first question is: **How does the spectrum change from one floor location to another?**

## 3. If you want a copy on your own computer

This is optional. `git clone` simply downloads a copy of the project. If Git or the command line is new to you, a collaborator can do this one-time setup with you.

```bash
git clone https://github.com/BlaineHeffron/HFIR_BG_Analysis.git
cd HFIR_BG_Analysis
./scripts/setup_analysis.sh --browser-only
./scripts/run_data_browser.sh
```

After the last command, open <http://localhost:8501> in a web browser. The setup prepares a local, read-only copy of the public catalog and spectra; it does not alter the published data.

## 4. What Figure 7 represents

Figure 7 uses the original downward-looking floor survey. The detector was measuring radiation arriving from the floor direction, so these spectra should not be described as a measurement of radiation arriving equally from every direction.

Later checks at some of the same locations are kept separate. They were collected at a different time and for a different purpose, so mixing them with the original survey would blur the comparison.

## 5. Reading and reusing the data thoughtfully

- The plotted spectra are detector readings (counts) as a function of gamma-ray energy. They are not, by themselves, an estimate of the incoming gamma-ray flux.
- The original fine-grained energy channels remain available. A grouped display version is provided only to make broad patterns easier to see.
- Measurements at high energy are sparse, so their uncertainty is naturally larger.
- If you plan to publish, share, or draw a scientific conclusion from a selection, keep the selected location, conditions, and source information with it.

## 6. Where to get help

Start with the browser and a simple question. For a guided first look, ask the presenter to walk through one location and one spectrum. The repository `README.md`, `docs/PUBLIC_ANALYSIS_GUIDE.md`, and `docs/FLOOR_SCAN_STATISTICS.md` contain the optional technical details for later.
"""
    md_path = output_dir / "HFIR_student_data_quickstart.md"
    md_path.write_text(markdown, encoding="utf-8")

    pages: list[Image.Image] = []
    sections = [
        (
            "HFIR data guide — begin with the question",
            [
                ("You can start here", "You do not need to know Git, Python, or the project folders to follow this guide. Begin with the browser and one question about the measurements."),
                ("What is a spectrum?", "A spectrum is a record of how many gamma rays the detector observed at each energy. It is the radiation fingerprint of one place and one time."),
                ("What is in the collection?", "There are 1,802 calibrated spectra from 354 measurement runs at 241 mapped locations. Each is linked to location, reactor condition, shielding, and measurement duration."),
                ("A useful first question", "How does the spectrum change from one floor location to another? Start there rather than trying to learn every file or term at once."),
            ],
        ),
        (
            "HFIR data guide — explore first, set up later",
            [
                ("Use the browser", "Ask your presenter or colleague for the browser link. Filter by location or condition, choose a measurement point, select one recording, and view its spectrum. You can change the displayed energy range or download the selected values."),
                ("Optional one-time setup", "git clone https://github.com/BlaineHeffron/HFIR_BG_Analysis.git\ncd HFIR_BG_Analysis\n./scripts/setup_analysis.sh --browser-only"),
                ("Start your local browser", "./scripts/run_data_browser.sh\n# then open http://localhost:8501"),
                ("What the setup does", "It downloads a local, read-only copy of the public catalog and spectra. If the command line is new to you, ask a collaborator to do this first setup with you."),
            ],
        ),
        (
            "HFIR data guide — use the measurements carefully",
            [
                ("Figure 7 in plain language", "Figure 7 uses the original downward-looking floor survey. It measures radiation arriving from the floor direction, not radiation arriving equally from every direction."),
                ("Keep like with like", "Later checks at some of the same locations are not folded into the original survey because they were collected at a different time and for a different purpose."),
                ("Detailed and easy-to-read views", "The original fine-grained energy channels are preserved. Grouped display values make broad patterns easier to see, especially at high energy where readings are sparse."),
                ("When you need help", "Start with one location and one spectrum. The README and the public-analysis guide explain the optional technical details when you are ready for them."),
            ],
        ),
    ]
    for page_number, (title, blocks) in enumerate(sections, start=1):
        page = Image.new("RGB", (1275, 1650), WHITE)
        draw = ImageDraw.Draw(page)
        draw.rectangle((0, 0, 1275, 125), fill=NAVY)
        draw.text((65, 38), title, font=font(34, bold=True), fill=WHITE)
        y = 165
        for heading, content in blocks:
            draw.text((70, y), heading, font=font(27, bold=True), fill=NAVY)
            y += 48
            is_code = any(token in content for token in ("./scripts/", ".venv/", "git clone"))
            if is_code:
                lines = content.count("\n") + 1
                height = 48 + lines * 33
                draw.rounded_rectangle((68, y, 1205, y + height), radius=12, fill="#17242e")
                draw.multiline_text((90, y + 20), content, font=font(19, mono=True), fill="#edf4f7", spacing=8)
                y += height + 38
            else:
                used = wrapped(draw, content, (75, y), 92, 23, color=TEXT, spacing=9)
                y += used + 42
        draw.text((70, 1595), f"Page {page_number}/3 | HFIR gamma-background data", font=font(17), fill=MUTED)
        page.save(output_dir / f"handout_page_{page_number}.png", optimize=True)
        pages.append(page)
    pdf_path = output_dir / "HFIR_student_data_quickstart.pdf"
    pages[0].save(pdf_path, "PDF", resolution=150, save_all=True, append_images=pages[1:])
    return md_path, pdf_path


def build_full_deck(output_dir: Path, browser_screenshot: Path) -> list[Image.Image]:
    public = REPO_ROOT / "analysis" / "public_analysis"
    report = REPO_ROOT / "reports" / "floor_scan_statistics"
    validation_path = output_dir / "e2e_validation.json"
    validation = json.loads(validation_path.read_text()) if validation_path.is_file() else None

    slide_specs: list[tuple[str, str | None]] = [
        ("Public data products and student workflow", None),
        ("Available analysis products", "From paper figures to reusable public data"),
        ("Public release at a glance", "A read-only analysis layer over the canonical release"),
        ("Unfolded gamma-flux spectra", "MIF reactor-on/off and Shield Center"),
        ("Shield configurations corresponding to Figure 14", "No added water/lead through seven layers plus floor lead"),
        ("Figure 7: recover the original scan population", "Avoid mixing later monitoring at the same coordinates"),
        ("What does an ‘average’ scan point look like?", "Routine acquisitions only: 100–400 s live time"),
        ("Representative point and suggested display binning", "HB4_DOWN_2: 216.25 s and 38,393 counts"),
        ("The point-to-point range is physically large", "Five spectra spanning total-count quantiles"),
        ("Binning recommendation", "Preserve native counts; provide a piecewise-binned convenience product"),
        ("What can be shared immediately", "Point manifest, all individual histograms, plots, and provenance"),
        ("Live browser demonstration", "Read-only exploration; no ROOT required"),
        ("Student quickstart", "Three commands to data exploration"),
        ("Student reproduction exercise", "Paper-facing products plus an independent native-channel export"),
        ("Paper reproducibility boundary", "Clear status for all 28 numbered figures"),
        ("Clean-clone validation", "The handout is tested, not aspirational"),
        ("Supplemental-release options", "Data products and metadata choices"),
    ]
    total = len(slide_specs)
    slides: list[Image.Image] = []

    slides.append(
        title_slide(
            "HFIR gamma-background\npublic data",
            "Paper-facing products, Figure 7 point statistics,\nand a student-ready analysis workflow",
            "Data and analysis overview | July 2026",
        )
    )

    image = base_slide(*slide_specs[1], 2, total)
    draw = ImageDraw.Draw(image)
    bullets(draw, [
        "Individual HPGe spectra underlie the downward-facing floor scan in Figure 7.",
        "Typical-point statistics guide the choice of supplemental histogram binning.",
        "Unfolded flux is available for MIF reactor-on/off and Shield Center, alongside Figure 14-like shield comparisons.",
        "The complete released data can be browsed, filtered, and exported by researchers and students.",
    ], (90, 165, 1470, 780), size=29)
    slides.append(image)

    image = base_slide(*slide_specs[2], 3, total)
    draw = ImageDraw.Draw(image)
    metric_card(draw, (80, 180, 430, 360), "354", "run records")
    metric_card(draw, (470, 180, 820, 360), "1,802", "calibrated spectrum files", GREEN)
    metric_card(draw, (860, 180, 1210, 360), "241", "mapped coordinate records", GOLD)
    metric_card(draw, (1250, 180, 1520, 360), "28", "paper figures inventoried", RED)
    bullets(draw, [
        "ROOT-free browser and CSV tools open the SQLite database read-only.",
        "Portable environment variables replace creator-machine paths; setup sanitizes the downloaded database copy.",
        "Official HFIR cycle dates are stored durably with day precision and source provenance.",
    ], (100, 450, 1500, 790), size=27)
    slides.append(image)

    image = base_slide(*slide_specs[3], 4, total)
    fit_image(image, public / "unfolded_key_locations.png", (60, 145, 1110, 825))
    draw = ImageDraw.Draw(image)
    bullets(draw, [
        "Official arXiv ancillary CSVs are the authoritative unfolded values.",
        "Isotropic and front-face responses bracket angular-model assumptions.",
        "This is a replot—not a new unfold—because response matrices and the external unfolder are not in the spectrum release.",
    ], (1140, 170, 1540, 790), size=21)
    slides.append(image)

    image = base_slide(*slide_specs[4], 5, total)
    fit_image(image, public / "shield_configuration_spectra.png", (80, 155, 1060, 810))
    draw = ImageDraw.Draw(image)
    bullets(draw, [
        "Figure 14 is recalculated from released raw spectra and the canonical database.",
        "Configurations: baseline; seven water layers; six layers plus floor lead; seven layers plus floor lead.",
        "CSV includes rates and statistical errors for independent reuse.",
    ], (1100, 180, 1530, 780), size=22)
    slides.append(image)

    image = base_slide(*slide_specs[5], 6, total)
    fit_image(image, report / "floor_scan_point_statistics.png", (50, 160, 1070, 815))
    draw = ImageDraw.Draw(image)
    bullets(draw, [
        "Figure 7 comes from the original downward-looking floor survey: 122 measurements at 117 locations.",
        "Most were routine measurements lasting about 2–7 minutes; those give us a fair picture of a typical survey point.",
        "Later checks used some of the same locations, but answered a different question—so they are kept separate here.",
        "Because the detector looked downward, these spectra describe radiation arriving from the floor direction, not from every direction around it.",
    ], (1100, 165, 1540, 800), size=20)
    slides.append(image)

    image = base_slide(*slide_specs[6], 7, total)
    draw = ImageDraw.Draw(image)
    metric_card(draw, (90, 180, 480, 385), "227.26 s", "median live time")
    metric_card(draw, (605, 180, 995, 385), "37,446", "median counts, 50–11,400 keV", GREEN)
    metric_card(draw, (1120, 180, 1510, 385), "178.09 s⁻¹", "median total count rate", GOLD)
    metric_card(draw, (90, 470, 670, 690), "4,290–223,525", "routine point count range", RED)
    metric_card(draw, (790, 470, 1510, 690), "HB4_DOWN_2", "robust representative: 216.25 s, 38,393 counts", BLUE)
    draw.text((90, 740), "‘Average’ is defined by robust distance to median live time, counts, and rate.", font=font(24), fill=MUTED)
    slides.append(image)

    image = base_slide(*slide_specs[7], 8, total)
    fit_image(image, report / "representative_point_adaptive_binning.png", (60, 140, 1540, 830))
    slides.append(image)

    image = base_slide(*slide_specs[8], 9, total)
    fit_image(image, report / "floor_scan_spectrum_quantiles.png", (60, 140, 1540, 830))
    slides.append(image)

    image = base_slide(*slide_specs[9], 10, total)
    draw = ImageDraw.Draw(image)
    simple_table(
        draw,
        ["Energy range", "Group size for display", "Typical readings per group", "Groups with a reading"],
        [
            ["50–1,000 keV", "2 keV", "46.0", "100%"],
            ["1,000–3,000 keV", "10 keV", "55.8", "100%"],
            ["3,000–7,000 keV", "50 keV", "48.3", "100%"],
            ["7,000–11,400 keV", "200 keV", "16.2", "77%"],
        ],
        (110, 180, 1490, 560),
        widths=[1.4, 1.0, 1.35, 1.2],
    )
    bullets(draw, [
        "Above 7 MeV, only 41% of 200-keV bins contain at least 10 counts at a typical point.",
        "Recommendation: distribute native calibrated channels plus this convenience binning and Poisson errors.",
    ], (130, 620, 1470, 810), size=25)
    slides.append(image)

    image = base_slide(*slide_specs[10], 11, total)
    draw = ImageDraw.Draw(image)
    bullets(draw, [
        "A point manifest with file/run IDs, location, live time, calibration, counts, and routine/extended flags.",
        "All 122 individual spectra in a 1.9 MB compressed CSV using the recommended piecewise binning.",
        "Native calibrated channels remain in the versioned public spectrum bundle.",
        "Representative and five-quantile plots, full bin-occupancy table, summary JSON, and exact regeneration command.",
        "Directionality and paper-reproducibility caveats travel with the package.",
    ], (100, 170, 1480, 790), size=29)
    slides.append(image)

    image = base_slide(*slide_specs[11], 12, total)
    draw = ImageDraw.Draw(image)
    if browser_screenshot.is_file():
        fit_image(image, browser_screenshot, (55, 145, 1120, 820))
    else:
        draw.rounded_rectangle((55, 145, 1120, 820), radius=18, fill=LIGHT, outline="#c9d2d8", width=2)
        wrapped(
            draw,
            "Optional live-demo screenshot not supplied. Start ./scripts/run_data_browser.sh, open localhost:8501, and pass a capture with --browser-screenshot.",
            (150, 380),
            55,
            27,
            color=MUTED,
            bold=True,
        )
    bullets(draw, [
        "Filter official cycle/state, shield, run text, and location.",
        "Select a map point, then a run and spectrum file.",
        "Change energy range, normalization, rebinning, log scale, and errors.",
        "Download the filtered run table or displayed calibrated bins.",
    ], (1150, 160, 1540, 800), size=20)
    slides.append(image)

    image = base_slide(*slide_specs[12], 13, total)
    draw = ImageDraw.Draw(image)
    code_box(draw, "git clone https://github.com/BlaineHeffron/HFIR_BG_Analysis.git\ncd HFIR_BG_Analysis\n./scripts/setup_analysis.sh --browser-only\n./scripts/run_data_browser.sh", (90, 180, 1510, 430), size=27)
    bullets(draw, [
        "Open http://localhost:8501.",
        "Search for position_scan_ or position_scan_3_HB4_DOWN_2.",
        "No CERN ROOT installation is needed for browsing, plotting, or CSV export.",
    ], (130, 500, 1460, 800), size=28)
    slides.append(image)

    image = base_slide(*slide_specs[13], 14, total)
    draw = ImageDraw.Draw(image)
    code_box(draw, ".venv/bin/python scripts/analyze_floor_scan_statistics.py\n\n.venv/bin/python scripts/export_public_data.py spectra \\\n  --file-id 1729 --normalization counts/s/keV \\\n  --emin 50 --emax 11400 \\\n  --output analysis/student/HB4_DOWN_2_native.csv", (80, 165, 1520, 510), size=23)
    code_box(draw, "# Linux x86-64 paper-facing products\n./scripts/setup_analysis.sh\nsource .env\n.venv/bin/python scripts/public_analysis.py all\n.venv/bin/python scripts/reproduce_paper.py --list", (80, 555, 1520, 800), size=22)
    slides.append(image)

    image = base_slide(*slide_specs[14], 15, total)
    draw = ImageDraw.Draw(image)
    metric_card(draw, (100, 180, 570, 390), "Figure 14", "recalculated from public measurements", GREEN)
    metric_card(draw, (650, 180, 1120, 390), "Figure 19", "three-location subset replotted from ancillary CSVs", BLUE)
    metric_card(draw, (1200, 180, 1510, 390), "26", "tracked; not recalculated at this checkpoint", GOLD)
    bullets(draw, [
        "Every numbered figure has a machine-readable status, inputs, publication artifact, and limitation.",
        "Figure 7 point spectra/statistics now have a supported workflow, while paper-exact contour styling remains legacy.",
        "External PROSPECT inputs and Geant4 response products are not silently treated as reproducible.",
    ], (120, 480, 1490, 790), size=27)
    slides.append(image)

    image = base_slide(*slide_specs[15], 16, total)
    draw = ImageDraw.Draw(image)
    if validation:
        metric_card(draw, (100, 185, 520, 390), validation.get("setup", "PASS"), "fresh browser-only and full setup", GREEN)
        metric_card(draw, (590, 185, 1010, 390), str(validation.get("spectrum_files", 1802)), "files found in clean clone", BLUE)
        metric_card(draw, (1080, 185, 1500, 390), validation.get("browser_health", "PASS"), "headless browser health check", GOLD)
        bullets(draw, validation.get("checks", []), (120, 490, 1490, 800), size=24)
    else:
        bullets(draw, ["Clean-clone validation will be recorded here after running the exact handout commands."], (120, 250, 1450, 700), size=32)
    slides.append(image)

    image = base_slide(*slide_specs[16], 17, total)
    draw = ImageDraw.Draw(image)
    bullets(draw, [
        "Share native channels only, or native channels plus the 2/10/50/200-keV convenience product?",
        "Include all 122 original acquisitions, or highlight the 104 routine four-minute points and flag extended follow-ups?",
        "Adopt explicit downward-collimated wording in the supplemental caption and metadata?",
        "Complete one final scan-point manifest and location review before release.",
        "Students can begin analysis immediately using the tested handout and browser.",
    ], (100, 170, 1490, 800), size=29)
    slides.append(image)
    return slides


def build_deck(output_dir: Path, browser_screenshot: Path) -> list[Image.Image]:
    """Build the concise Figure 7 and public-data presentation."""
    del output_dir
    report = REPO_ROOT / "reports" / "floor_scan_statistics"
    slide_specs: list[tuple[str, str | None]] = [
        ("HFIR gamma background: Figure 7", None),
        ("Dataset at a glance", "Released calibrated spectra and location metadata"),
        ("Figure 7 selection", "Original downward-looking floor scan"),
        ("Representative spectrum", "Acquisition nearest the floor-scan median"),
        ("Native and display spectra", "Representative point: fine channels and binned view"),
        ("Adaptive binning", "Native channels retained; bins are a viewing convenience"),
        ("Map and browser", "Detector orientation and spectrum access"),
        ("Analysis workflow", "Point → acquisition → spectrum"),
    ]
    total = len(slide_specs)
    slides: list[Image.Image] = [
        title_slide(
            "HFIR gamma background\nfor Figure 7",
            "Floor-scan data, spectrum products,\nand public analysis tools",
            "HFIR gamma-background data | July 2026",
        )
    ]

    image = base_slide(*slide_specs[1], 2, total)
    draw = ImageDraw.Draw(image)
    metric_card(draw, (100, 205, 525, 445), "354", "runs")
    metric_card(draw, (585, 205, 1010, 445), "1,802", "calibrated spectra", GREEN)
    metric_card(draw, (1070, 205, 1495, 445), "241", "mapped locations", GOLD)
    draw.text((100, 570), "Energy histograms + run, location, shielding, and operating-state metadata", font=font(30), fill=TEXT)
    slides.append(image)

    image = base_slide(*slide_specs[2], 3, total)
    fit_image(image, report / "floor_scan_point_statistics.png", (50, 145, 1120, 820))
    draw = ImageDraw.Draw(image)
    metric_card(draw, (1160, 185, 1515, 355), "122", "original spectra", GOLD)
    metric_card(draw, (1160, 395, 1515, 565), "117", "locations", GREEN)
    draw.text((1160, 640), "Detector axis", font=font(24), fill=MUTED)
    draw.text((1160, 675), "↓  downward", font=font(34, bold=True), fill=NAVY)
    draw.text((1160, 760), "Later follow-ups excluded", font=font(20), fill=MUTED)
    slides.append(image)

    image = base_slide(*slide_specs[3], 4, total)
    draw = ImageDraw.Draw(image)
    metric_card(draw, (90, 180, 480, 385), "about 4 min", "typical measurement length")
    metric_card(draw, (605, 180, 995, 385), "37,446", "typical detector readings", GREEN)
    metric_card(draw, (1120, 180, 1510, 385), "178 / sec", "typical reading rate", GOLD)
    metric_card(draw, (90, 470, 670, 690), "4,290–223,525", "range across routine locations", RED)
    metric_card(draw, (790, 470, 1510, 690), "Representative", "an example near the middle: 216 s and 38,393 readings", BLUE)
    draw.text((90, 760), "Selection metric: nearest median in (acquisition time, total counts, count rate)", font=font(23), fill=MUTED)
    slides.append(image)

    image = base_slide(*slide_specs[4], 5, total)
    fit_image(image, report / "representative_point_adaptive_binning.png", (60, 140, 1540, 830))
    slides.append(image)

    image = base_slide(*slide_specs[5], 6, total)
    draw = ImageDraw.Draw(image)
    simple_table(
        draw,
        ["Energy range", "Group size for display", "Typical readings per group", "Groups with a reading"],
        [
            ["50–1,000 keV", "2 keV", "46.0", "100%"],
            ["1,000–3,000 keV", "10 keV", "55.8", "100%"],
            ["3,000–7,000 keV", "50 keV", "48.3", "100%"],
            ["7,000–11,400 keV", "200 keV", "16.2", "77%"],
        ],
        (110, 180, 1490, 560),
        widths=[1.4, 1.0, 1.35, 1.2],
    )
    draw.text((130, 665), "Adaptive bins improve visual comparison; the released product preserves native-channel counts.", font=font(27), fill=TEXT)
    slides.append(image)

    image = base_slide(*slide_specs[6], 7, total)
    draw = ImageDraw.Draw(image)
    if browser_screenshot.is_file():
        fit_image(
            image,
            browser_screenshot,
            (55, 175, 1130, 720),
            crop_fraction=(0.1875, 0.37, 0.875, 0.84),
        )
    else:
        draw.rounded_rectangle((55, 145, 1120, 820), radius=18, fill=LIGHT, outline="#c9d2d8", width=2)
        wrapped(
            draw,
            "This is where the presenter can demonstrate selecting a location, opening one spectrum, and changing the display.",
            (150, 390), 55, 27, color=MUTED, bold=True,
        )
    draw.text((1160, 180), "Map encoding", font=font(28, bold=True), fill=NAVY)
    draw.text((1160, 250), "★", font=font(44), fill=GOLD)
    draw.text((1210, 260), "detector axis down", font=font(25), fill=TEXT)
    draw.text((1160, 350), "→", font=font(46, bold=True), fill=RED)
    draw.text((1210, 360), "tilted: horizontal projection", font=font(25), fill=TEXT)
    draw.text((1160, 500), "Select marker", font=font(26, bold=True), fill=NAVY)
    draw.text((1160, 545), "→ select acquisition", font=font(24), fill=TEXT)
    draw.text((1160, 585), "→ plot or export spectrum", font=font(24), fill=TEXT)
    slides.append(image)

    image = base_slide(*slide_specs[7], 8, total)
    draw = ImageDraw.Draw(image)
    metric_card(draw, (80, 215, 510, 490), "1. Locate", "map marker or metadata filter")
    metric_card(draw, (590, 215, 1020, 490), "2. Select", "acquisition at a point", GREEN)
    metric_card(draw, (1100, 215, 1530, 490), "3. Inspect", "spectrum or CSV export", GOLD)
    draw.text((110, 625), "Read-only browser; analysis starts from the native calibrated spectrum.", font=font(30), fill=TEXT)
    slides.append(image)
    return slides


def main() -> int:
    args = parse_args()
    output_dir = args.output_dir.expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    required = [
        REPO_ROOT / "reports" / "floor_scan_statistics" / "floor_scan_point_statistics.png",
        REPO_ROOT / "reports" / "floor_scan_statistics" / "representative_point_adaptive_binning.png",
    ]
    missing = [str(path) for path in required if not path.is_file()]
    if missing:
        raise SystemExit("Missing presentation assets:\n" + "\n".join(missing))

    slides = build_deck(output_dir, args.browser_screenshot.expanduser().resolve())
    pptx_path, preview_path = save_slides(slides, output_dir)
    md_path, handout_path = write_handout(output_dir)
    print(f"Wrote {len(slides)} slides to {pptx_path}")
    print(f"Wrote slide preview to {preview_path}")
    print(f"Wrote student handout to {md_path} and {handout_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
